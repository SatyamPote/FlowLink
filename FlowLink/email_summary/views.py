import os
import shutil
import imaplib
import email
from email.header import decode_header
import pandas as pd
import json
from django.shortcuts import render, redirect
from django.contrib.auth import authenticate, login as auth_login, logout as auth_logout
from django.contrib.auth.forms import AuthenticationForm, UserCreationForm
from django.conf import settings
from django.http import HttpResponse
from pathlib import Path
from docx import Document
from pptx import Presentation
from fpdf import FPDF
import img2pdf
from PIL import Image
from PyPDF2 import PdfReader
from openpyxl import load_workbook
from reportlab.pdfgen import canvas
import io

def login_view(request):
    if request.method == 'POST':
        form = AuthenticationForm(request, data=request.POST)
        if form.is_valid():
            username = form.cleaned_data.get('username')
            password = form.cleaned_data.get('password')
            user = authenticate(username=username, password=password)
            if user is not None:
                auth_login(request, user)
                return redirect('dashboard')
    else:
        form = AuthenticationForm()
    return render(request, 'login.html', {'form': form})

def register_view(request):
    if request.method == 'POST':
        form = UserCreationForm(request.POST)
        if form.is_valid():
            user = form.save()
            auth_login(request, user)
            return redirect('dashboard')
    else:
        form = UserCreationForm()
    return render(request, 'register.html', {'form': form})

def dashboard(request):
    return render(request, 'dashboard.html')

def logout_view(request):
    auth_logout(request)
    return redirect('login')

def excel_view(request):
    emails = []
    data_extracted = False
    if request.method == 'POST':
        mail = imaplib.IMAP4_SSL(settings.EMAIL_HOST)
        mail.login(settings.EMAIL_HOST_USER, settings.EMAIL_HOST_PASSWORD)
        mail.select("inbox")

        result, data = mail.search(None, "ALL")
        email_ids = data[0].split()[-100:]

        emails_list = []

        for eid in email_ids:
            result, msg_data = mail.fetch(eid, "(RFC822)")
            raw_email = msg_data[0][1]
            msg = email.message_from_bytes(raw_email)
            
            subject, encoding = decode_header(msg["Subject"])[0]
            if isinstance(subject, bytes):
                subject = subject.decode(encoding if encoding else "utf-8")
            
            from_ = msg.get("From")
            date_ = msg.get("Date")
            
            emails_list.append({"subject": subject, "from": from_, "date": date_})

        df = pd.DataFrame(emails_list)
        json_file_path = Path(settings.BASE_DIR) / 'emails_data.json'
        df.to_json(json_file_path, orient='records')

        request.session['emails_data_path'] = str(json_file_path)
        emails = emails_list
        data_extracted = True

    return render(request, 'emails/excel.html', {'emails': emails, 'data_extracted': data_extracted})

def download_json(request):
    json_file_path = request.session.get('emails_data_path')
    if not json_file_path:
        return redirect('excel_view')

    with open(json_file_path, 'r') as f:
        data = f.read()
    response = HttpResponse(data, content_type='application/json')
    response['Content-Disposition'] = 'attachment; filename="emails_data.json"'
    return response

def download_csv(request):
    json_file_path = request.session.get('emails_data_path')
    if not json_file_path:
        return redirect('excel_view')

    with open(json_file_path, 'r') as f:
        emails = json.load(f)

    df = pd.DataFrame(emails)
    response = HttpResponse(content_type='text/csv')
    response['Content-Disposition'] = 'attachment; filename="emails_data.csv"'
    df.to_csv(path_or_buf=response, index=False)
    
    return response

def file_organizer_view(request):
    if request.method == 'POST':
        directory = request.POST.get('directory')
        if not os.path.isdir(directory):
            return render(request, 'file_organizer.html', {'error': 'Invalid directory'})

        extensions = {
            'Documents': ['.pdf', '.doc', '.docx', '.txt'],
            'Images': ['.jpg', '.jpeg', '.png', '.gif', '.bmp'],
            'Audio': ['.mp3', '.wav', '.aac'],
            'Video': ['.mp4', '.avi', '.mkv'],
            'Archives': ['.zip', '.rar', '.tar'],
            'Others': []
        }

        for filename in os.listdir(directory):
            file_path = os.path.join(directory, filename)
            if os.path.isfile(file_path):
                file_ext = os.path.splitext(filename)[1].lower()
                moved = False
                for folder, exts in extensions.items():
                    if file_ext in exts:
                        folder_path = os.path.join(directory, folder)
                        if not os.path.exists(folder_path):
                            os.makedirs(folder_path)
                        shutil.move(file_path, os.path.join(folder_path, filename))
                        moved = True
                        break
                if not moved:
                    other_folder_path = os.path.join(directory, 'Others')
                    if not os.path.exists(other_folder_path):
                        os.makedirs(other_folder_path)
                    shutil.move(file_path, os.path.join(other_folder_path, filename))

        return render(request, 'file_organizer.html', {'success': 'Files organized successfully!'})

    return render(request, 'file_organizer.html')

# File Conversion Functions
def jpg_to_pdf(file_path, output_buffer):
    image = Image.open(file_path)
    pdf_bytes = img2pdf.convert(image.filename)
    output_buffer.write(pdf_bytes)

def pdf_to_jpg(pdf_bytes, output_buffer):
    from pdf2image import convert_from_bytes
    images = convert_from_bytes(pdf_bytes)
    for i, image in enumerate(images):
        image.save(output_buffer, format='JPEG')

def word_to_pdf(doc_bytes, output_buffer):
    doc = Document(io.BytesIO(doc_bytes))
    pdf = canvas.Canvas(output_buffer)
    for paragraph in doc.paragraphs:
        pdf.drawString(100, 750, paragraph.text)
        pdf.showPage()
    pdf.save()

def pdf_to_word(pdf_bytes, output_buffer):
    reader = PdfReader(io.BytesIO(pdf_bytes))
    doc = Document()
    for page in reader.pages:
        doc.add_paragraph(page.extract_text())
    output_stream = io.BytesIO()
    doc.save(output_stream)
    output_buffer.write(output_stream.getvalue())

def ppt_to_pdf(ppt_bytes, output_buffer):
    prs = Presentation(io.BytesIO(ppt_bytes))
    pdf = canvas.Canvas(output_buffer)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                pdf.drawString(100, 750, shape.text)
                pdf.showPage()
    pdf.save()

def excel_to_pdf(excel_bytes, output_buffer):
    wb = load_workbook(filename=io.BytesIO(excel_bytes))
    sheet = wb.active
    pdf = canvas.Canvas(output_buffer)
    for row in sheet.iter_rows(values_only=True):
        pdf.drawString(100, 750, " ".join([str(cell) for cell in row]))
        pdf.showPage()
    pdf.save()

def html_to_pdf(html_string, output_buffer):
    from weasyprint import HTML
    HTML(string=html_string).write_pdf(output_buffer)

# File Converter View
def convert_file(file_bytes, output_format, new_file_name, file_ext):
    output_buffer = io.BytesIO()

    try:
        if output_format == 'pdf' and file_ext == '.jpg':
            jpg_to_pdf(io.BytesIO(file_bytes), output_buffer)
        elif output_format == 'jpg' and file_ext == '.pdf':
            pdf_to_jpg(file_bytes, output_buffer)
        elif output_format == 'pdf' and file_ext == '.docx':
            word_to_pdf(file_bytes, output_buffer)
        elif output_format == 'docx' and file_ext == '.pdf':
            pdf_to_word(file_bytes, output_buffer)
        elif output_format == 'pdf' and file_ext == '.pptx':
            ppt_to_pdf(file_bytes, output_buffer)
        elif output_format == 'pdf' and file_ext == '.xlsx':
            excel_to_pdf(file_bytes, output_buffer)
        elif output_format == 'pdf' and file_ext == '.html':
            html_to_pdf(file_bytes.decode('utf-8'), output_buffer)
        else:
            return None
        return output_buffer.getvalue()
    except Exception as e:
        print(f"Error converting file: {e}")
        return None

def file_converter_view(request):
    if request.method == 'POST':
        if 'file' not in request.FILES:
            return render(request, 'file_converter.html', {'error': 'No file uploaded! Please upload a file.'})
        uploaded_file = request.FILES['file']
        output_format = request.POST.get('output_format')
        new_file_name = request.POST.get('new_file_name', '').strip()
        file_ext = os.path.splitext(uploaded_file.name)[1].lower()
        file_bytes = uploaded_file.read()

        output_bytes = convert_file(file_bytes, output_format, new_file_name, file_ext)
        if output_bytes:
            response = HttpResponse(output_bytes, content_type='application/octet-stream')
            response['Content-Disposition'] = f'attachment; filename="{new_file_name or uploaded_file.name}.{output_format}"'
            return response
        else:
            return render(request, 'file_converter.html', {'error': 'Conversion failed! Please provide a proper file or try a different format.'})
    return render(request, 'file_converter.html')

        